import argparse
import os as os
import pandas as pd
import numpy as np
from venn import venn
from matplotlib import pyplot as plt

from parsers import athlatesclassIIparser
from parsers import hlaminerclassIIparser
from parsers import optitypeclassIparser
from parsers import seq2hlaclassIIparser
from parsers import athlatesclassIparser
from parsers import hlaminerclassIparser
from parsers import polysolverclassIparser
from parsers import seq2hlaclassIparser
from dframes import dataframes

from hlatools import hlatoolsAthlatesClassIINormal
from hlatools import hlatoolsAthlatesClassIIRNA
from hlatools import hlatoolsAthlatesClassIITumor
from hlatools import hlatoolsAthlatesClassINormal
from hlatools import hlatoolsAthlatesClassIRNA
from hlatools import hlatoolsAthlatesClassITumor
from hlatools import hlatoolsHLAminerClassIIRNA
from hlatools import hlatoolsHLAminerClassIRNA
from hlatools import hlatoolsOptitypeClassINormal
from hlatools import hlatoolsOptitypeClassIRNA
from hlatools import hlatoolsOptitypeClassITumor
from hlatools import hlatoolsPolysolverClassINormal
from hlatools import hlatoolsPolysolverClassIRNA
from hlatools import hlatoolsPolysolverClassITumor
from hlatools import hlatoolsSeq2HLAClassIINormal
from hlatools import hlatoolsSeq2HLAClassIIRNA
from hlatools import hlatoolsSeq2HLAClassIITumor
from hlatools import hlatoolsSeq2HLAClassINormal
from hlatools import hlatoolsSeq2HLAClassIRNA
from hlatools import hlatoolsSeq2HLAClassITumor

   
from consensus import consensusAllClassI
from consensus import consensusAllClassII
from consensus import consensusAthlatesClassI
from consensus import consensusAthlatesClassII
from consensus import consensusOptitypeClassI
from consensus import consensusPolysolverClassI
from consensus import consensusSeq2HLAClassI
from consensus import consensusSeq2HLAClassII


parser = argparse.ArgumentParser(description = "help message")

parser.add_argument("-i", "--InputDir", help = "Input directory")
parser.add_argument("-s", "--SampleID", help = "Sample ID")

args = parser.parse_args()

SampleID = args.SampleID
CurrDir = args.InputDir
CWD = os.getcwd()

OutDir = CurrDir # dir as arguement # to use for output csv plots and pptx
AllHLATools = os.listdir(os.path.join(os.getcwd(), CurrDir))
AllHLATools



AthlatesClassI=dataframes.dataframes()[0]
AthlatesClassII=dataframes.dataframes()[1]
OptiTypeNormalTumorRNA=dataframes.dataframes()[2]
PolysolverNormalTumorRNA=dataframes.dataframes()[3]
Seq2HLAClassINormalTumorRNA=dataframes.dataframes()[4]
Seq2HLAClassIINormalTumorRNA=dataframes.dataframes()[5]
HLAMinerClassIRNATumorHLA=dataframes.dataframes()[6]
HLAMinerClassIIRNATumorHLA=dataframes.dataframes()[7]

                            
                        
                        

print("AllHLATools are: ", AllHLATools)                
#get all HLA calls in dataframes
AthlatesClassI=hlatoolsAthlatesClassINormal.getHLA(CWD, CurrDir, AllHLATools, AthlatesClassI)[0]
AthlatesClassI=hlatoolsAthlatesClassIRNA.getHLA(CWD, CurrDir, AllHLATools, AthlatesClassI)[0]
AthlatesClassI=hlatoolsAthlatesClassITumor.getHLA(CWD, CurrDir, AllHLATools, AthlatesClassI)[0]

AthlatesClassII=hlatoolsAthlatesClassIINormal.getHLA(CWD, CurrDir, AllHLATools, AthlatesClassII)[0]
AthlatesClassII=hlatoolsAthlatesClassIIRNA.getHLA(CWD, CurrDir, AllHLATools, AthlatesClassII)[0]
AthlatesClassII=hlatoolsAthlatesClassIITumor.getHLA(CWD, CurrDir, AllHLATools, AthlatesClassII)[0]

OptiTypeNormalTumorRNA=hlatoolsOptitypeClassINormal.getHLA(CWD, CurrDir, AllHLATools, OptiTypeNormalTumorRNA)[0]
OptiTypeNormalTumorRNA=hlatoolsOptitypeClassITumor.getHLA(CWD, CurrDir, AllHLATools, OptiTypeNormalTumorRNA)[0]
OptiTypeNormalTumorRNA=hlatoolsOptitypeClassIRNA.getHLA(CWD, CurrDir, AllHLATools, OptiTypeNormalTumorRNA)[0]

PolysolverNormalTumorRNA=hlatoolsPolysolverClassINormal.getHLA(CWD, CurrDir, AllHLATools, PolysolverNormalTumorRNA)[0]
PolysolverNormalTumorRNA=hlatoolsPolysolverClassIRNA.getHLA(CWD, CurrDir, AllHLATools, PolysolverNormalTumorRNA)[0]
PolysolverNormalTumorRNA=hlatoolsPolysolverClassITumor.getHLA(CWD, CurrDir, AllHLATools, PolysolverNormalTumorRNA)[0]

Seq2HLAClassINormalTumorRNA=hlatoolsSeq2HLAClassINormal.getHLA(CWD, CurrDir, AllHLATools, Seq2HLAClassINormalTumorRNA)[0]
Seq2HLAClassINormalTumorRNA=hlatoolsSeq2HLAClassIRNA.getHLA(CWD, CurrDir, AllHLATools, Seq2HLAClassINormalTumorRNA)[0]
Seq2HLAClassINormalTumorRNA=hlatoolsSeq2HLAClassITumor.getHLA(CWD, CurrDir, AllHLATools, Seq2HLAClassINormalTumorRNA)[0]

Seq2HLAClassIINormalTumorRNA=hlatoolsSeq2HLAClassIINormal.getHLA(CWD, CurrDir, AllHLATools, Seq2HLAClassIINormalTumorRNA)[0]
Seq2HLAClassIINormalTumorRNA=hlatoolsSeq2HLAClassIIRNA.getHLA(CWD, CurrDir, AllHLATools, Seq2HLAClassIINormalTumorRNA)[0]
Seq2HLAClassIINormalTumorRNA=hlatoolsSeq2HLAClassIITumor.getHLA(CWD, CurrDir, AllHLATools, Seq2HLAClassIINormalTumorRNA)[0]

HLAMinerClassIRNATumorHLA=hlatoolsHLAminerClassIRNA.getHLA(CWD, CurrDir, AllHLATools, HLAMinerClassIRNATumorHLA)[0]
HLAMinerClassIIRNATumorHLA=hlatoolsHLAminerClassIIRNA.getHLA(CWD, CurrDir, AllHLATools, HLAMinerClassIIRNATumorHLA)[0]



print("AthlatesClassI:")
print(AthlatesClassI)
print("Sorted AthlatesClassI:")
print(pd.DataFrame({"DNA-Normal": pd.Series(sorted(AthlatesClassI["DNA-Normal"])), "DNA-Tumor": pd.Series(sorted(AthlatesClassI["DNA-Tumor"])), "RNA-Tumor": pd.Series(sorted(AthlatesClassI["RNA-Tumor"]))}))
print("AthlatesClassII:")
print(AthlatesClassII)
print("Sorted AthlatesClassII:")
print(pd.DataFrame({"DNA-Normal": pd.Series(sorted(AthlatesClassII["DNA-Normal"])), "DNA-Tumor": pd.Series(sorted(AthlatesClassII["DNA-Tumor"])), "RNA-Tumor": pd.Series(sorted(AthlatesClassII["RNA-Tumor"]))}))
print("OptiTypeNormalTumorRNA:")
print(OptiTypeNormalTumorRNA)
print("Sorted OptiTypeNormalTumorRNA:")
print(pd.DataFrame({"DNA-Normal": pd.Series(sorted(OptiTypeNormalTumorRNA["DNA-Normal"])), "DNA-Tumor": pd.Series(sorted(OptiTypeNormalTumorRNA["DNA-Tumor"])), "RNA-Tumor": pd.Series(sorted(OptiTypeNormalTumorRNA["RNA-Tumor"]))}))
print("PolysolverNormalTumorRNA:")
print(PolysolverNormalTumorRNA)
print("Sorted PolysolverNormalTumorRNA:")
print(pd.DataFrame({"DNA-Normal": pd.Series(sorted(PolysolverNormalTumorRNA["DNA-Normal"])), "DNA-Tumor": pd.Series(sorted(PolysolverNormalTumorRNA["DNA-Tumor"])), "RNA-Tumor": pd.Series(sorted(PolysolverNormalTumorRNA["RNA-Tumor"]))}))
print("Seq2HLAClassINormalTumorRNA:")
print(Seq2HLAClassINormalTumorRNA)
print("Sorted Seq2HLAClassINormalTumorRNA:")
print(pd.DataFrame({"DNA-Normal": pd.Series(sorted(Seq2HLAClassINormalTumorRNA["DNA-Normal"])), "DNA-Tumor": pd.Series(sorted(Seq2HLAClassINormalTumorRNA["DNA-Tumor"])), "RNA-Tumor": pd.Series(sorted(Seq2HLAClassINormalTumorRNA["RNA-Tumor"]))}))
print("Seq2HLAClassIINormalTumorRNA:")
print(Seq2HLAClassIINormalTumorRNA)
print("Sorted Seq2HLAClassIINormalTumorRNA:")
print(pd.DataFrame({"DNA-Normal": pd.Series(sorted(Seq2HLAClassIINormalTumorRNA["DNA-Normal"])), "DNA-Tumor": pd.Series(sorted(Seq2HLAClassIINormalTumorRNA["DNA-Tumor"])), "RNA-Tumor": pd.Series(sorted(Seq2HLAClassIINormalTumorRNA["RNA-Tumor"]))}))
print("HLAMinerClassIRNATumorHLA:")
print(HLAMinerClassIRNATumorHLA)
print("HLAMinerClassIIRNATumorHLA:")
print(HLAMinerClassIIRNATumorHLA)


AthlatesClassI_Consensus=consensusAthlatesClassI.getconsensusAthlatesClassI(AthlatesClassI)[0]
AthlatesClassII_Consensus=consensusAthlatesClassII.getconsensusAthlatesClassII(AthlatesClassII)[0]
OptiTypeNormalTumorRNA_Consensus=consensusOptitypeClassI.getconsensusOptitypeClassI(OptiTypeNormalTumorRNA)[0]
PolysolverNormalTumorRNA_Consensus=consensusPolysolverClassI.getconsensusPolysolverClassI(PolysolverNormalTumorRNA)[0]
Seq2HLAClassINormalTumorRNA_Consensus=consensusSeq2HLAClassI.getconsensusSeq2HLAClassI(Seq2HLAClassINormalTumorRNA)[0]
Seq2HLAClassIINormalTumorRNA_Consensus=consensusSeq2HLAClassII.getconsensusSeq2HLAClassII(Seq2HLAClassIINormalTumorRNA)[0]
#ClassI_Consensus=consensusAllClassI.
#ClassII_Consensus=consensusAllClassII.

    
#Final consensus Class I & II HLA list    
ConsensusClassIHLA = consensusAllClassI.FindConsensusClassI(AthlatesClassI_Consensus, OptiTypeNormalTumorRNA_Consensus, PolysolverNormalTumorRNA_Consensus, Seq2HLAClassINormalTumorRNA_Consensus, HLAMinerClassIRNATumorHLA)
print("@@@@@@@@@@@@@@@@ ConsensusClassIHLA: ", ConsensusClassIHLA)
#print("@@@@@@@@@@@@@@@@ ConsensusClassIHLA: ", sorted(ConsensusClassIHLA))

ConsensusClassIIHLA = consensusAllClassII.FindConsensusClassII(AthlatesClassII_Consensus, Seq2HLAClassIINormalTumorRNA_Consensus, HLAMinerClassIIRNATumorHLA)
print("@@@@@@@@@@@@@@@@ ConsensusClassIIHLA: ", sorted(ConsensusClassIIHLA))
print(type(ConsensusClassIIHLA))

#Final table for HLA list
#print("AthlatesClassI:")
#print(AthlatesClassI)
print("Sorted AthlatesClassI:")
AthlatesClassI = pd.DataFrame({"DNA-Normal": pd.Series(sorted(AthlatesClassI["DNA-Normal"])), "DNA-Tumor": pd.Series(sorted(AthlatesClassI["DNA-Tumor"])), "RNA-Tumor": pd.Series(sorted(AthlatesClassI["RNA-Tumor"])), "Consensus": pd.Series(sorted(AthlatesClassI_Consensus))})
print(AthlatesClassI)
#print("AthlatesClassII:")
#print(AthlatesClassII)
print("Sorted AthlatesClassII:")
AthlatesClassII = pd.DataFrame({"DNA-Normal": pd.Series(sorted(AthlatesClassII["DNA-Normal"])), "DNA-Tumor": pd.Series(sorted(AthlatesClassII["DNA-Tumor"])), "RNA-Tumor": pd.Series(sorted(AthlatesClassII["RNA-Tumor"])), "Consensus": pd.Series(sorted(AthlatesClassII_Consensus))})
print(AthlatesClassII)
#print("OptiTypeNormalTumorRNA:")
#print(OptiTypeNormalTumorRNA)
print("Sorted OptiTypeNormalTumorRNA:")
OptiTypeNormalTumorRNA = pd.DataFrame({"DNA-Normal": pd.Series(sorted(OptiTypeNormalTumorRNA["DNA-Normal"])), "DNA-Tumor": pd.Series(sorted(OptiTypeNormalTumorRNA["DNA-Tumor"])), "RNA-Tumor": pd.Series(sorted(OptiTypeNormalTumorRNA["RNA-Tumor"])), "Consensus": pd.Series(sorted(OptiTypeNormalTumorRNA_Consensus))})
print(OptiTypeNormalTumorRNA)
#print("PolysolverNormalTumorRNA:")
#print(PolysolverNormalTumorRNA)
print("Sorted PolysolverNormalTumorRNA:")
PolysolverNormalTumorRNA = pd.DataFrame({"DNA-Normal": pd.Series(sorted(PolysolverNormalTumorRNA["DNA-Normal"])), "DNA-Tumor": pd.Series(sorted(PolysolverNormalTumorRNA["DNA-Tumor"])), "RNA-Tumor": pd.Series(sorted(PolysolverNormalTumorRNA["RNA-Tumor"])), "Consensus": pd.Series(sorted(PolysolverNormalTumorRNA_Consensus))})
print(PolysolverNormalTumorRNA)
#print("Seq2HLAClassINormalTumorRNA:")
#print(Seq2HLAClassINormalTumorRNA)
print("Sorted Seq2HLAClassINormalTumorRNA:")
Seq2HLAClassINormalTumorRNA = pd.DataFrame({"DNA-Normal": pd.Series(sorted(Seq2HLAClassINormalTumorRNA["DNA-Normal"])), "DNA-Tumor": pd.Series(sorted(Seq2HLAClassINormalTumorRNA["DNA-Tumor"])), "RNA-Tumor": pd.Series(sorted(Seq2HLAClassINormalTumorRNA["RNA-Tumor"])), "Consensus": pd.Series(sorted(Seq2HLAClassINormalTumorRNA_Consensus))})
print(Seq2HLAClassINormalTumorRNA)
#print("Seq2HLAClassIINormalTumorRNA:")
#print(Seq2HLAClassIINormalTumorRNA)
print("Sorted Seq2HLAClassIINormalTumorRNA:")
Seq2HLAClassIINormalTumorRNA = pd.DataFrame({"DNA-Normal": pd.Series(sorted(Seq2HLAClassIINormalTumorRNA["DNA-Normal"])), "DNA-Tumor": pd.Series(sorted(Seq2HLAClassIINormalTumorRNA["DNA-Tumor"])), "RNA-Tumor": pd.Series(sorted(Seq2HLAClassIINormalTumorRNA["RNA-Tumor"])), "Consensus": pd.Series(sorted(Seq2HLAClassIINormalTumorRNA_Consensus))})
print(Seq2HLAClassIINormalTumorRNA)
print("HLAMinerClassIRNATumorHLA:")
print(HLAMinerClassIRNATumorHLA)
print("HLAMinerClassIIRNATumorHLA:")
print(HLAMinerClassIIRNATumorHLA)

#All tools Class I consensus HLA list table
AllClassIConsensusTable = pd.DataFrame({"Athlates": pd.Series(sorted(AthlatesClassI_Consensus)), "OptiType": pd.Series(sorted(OptiTypeNormalTumorRNA_Consensus)), "Polysolver": pd.Series(sorted(PolysolverNormalTumorRNA_Consensus)), "Seq2HLA": pd.Series(sorted(Seq2HLAClassINormalTumorRNA_Consensus)), "HLAMiner": pd.Series(sorted(HLAMinerClassIRNATumorHLA)), "Consensus": pd.Series(list(ConsensusClassIHLA))})
print("AllClassIConsensusTable")
print(AllClassIConsensusTable)
#All tools Class II consensus HLA list table
AllClassIIConsensusTable = pd.DataFrame({"Athlates": pd.Series(sorted(AthlatesClassII_Consensus)), "Seq2HLA": pd.Series(sorted(Seq2HLAClassIINormalTumorRNA_Consensus)), "HLAMiner": pd.Series(sorted(HLAMinerClassIIRNATumorHLA)), "Consensus": pd.Series(sorted(ConsensusClassIIHLA))})
print("AllClassIIConsensusTable")
print(AllClassIIConsensusTable)


#Save all tables to csv
AthlatesClassI.to_csv(OutDir+ "/" + SampleID + "-Athlates-Class-I-DNA-Normal-DNA-Tumor-RNA-Tumor.csv", sep=',', index=False)

OptiTypeNormalTumorRNA.to_csv(OutDir+ "/" + SampleID + "-Optitype-DNA-Normal-DNA-Tumor-RNA-Tumor.csv", sep=',', index=False)

PolysolverNormalTumorRNA.to_csv(OutDir+ "/" + SampleID + "-Polysolver-DNA-Normal-DNA-Tumor-RNA-Tumor.csv", sep=',', index=False)

Seq2HLAClassINormalTumorRNA.to_csv(OutDir+ "/" + SampleID + "-Seq2HLA-Class-I-DNA-Normal-DNA-Tumor-RNA-Tumor.csv", sep=',', index=False)


AthlatesClassII.to_csv(OutDir+ "/" + SampleID + "-Athlates-Class-II-DNA-Normal-DNA-Tumor-RNA-Tumor.csv", sep=',', index=False)

Seq2HLAClassIINormalTumorRNA.to_csv(OutDir+ "/" + SampleID + "-Seq2HLA-Class-II-DNA-Normal-DNA-Tumor-RNA-Tumor.csv", sep=',', index=False)

AllClassIConsensusTable.to_csv(OutDir+ "/" + SampleID + "-HLA-List-Class-I-Only.csv", sep=',', index=False)
AllClassIIConsensusTable.to_csv(OutDir+ "/" + SampleID + "-HLA-List-Class-II-Only.csv", sep=',', index=False)

print(OutDir)
print(SampleID)


#Class I HLA Only
HLAListClassI = pd.read_csv(OutDir+ "/" + SampleID + "-HLA-List-Class-I-Only.csv", header=0,sep=',')
HLAMinerClassI = set(item for item in HLAListClassI['HLAMiner'] if not(pd.isnull(item))==True)
#ClinicalClassI = set(item for item in HLAListClassI['Clinical-HLA'] if not(pd.isnull(item))==True)
seq2HLAClassI = set(item for item in HLAListClassI['Seq2HLA'] if not(pd.isnull(item))==True)
PolysolverClassI = set(item for item in HLAListClassI['Polysolver'] if not(pd.isnull(item))==True)
OptitypeClassI = set(item for item in HLAListClassI['OptiType'] if not(pd.isnull(item))==True)
AthlatesClassI = set(item for item in HLAListClassI['Athlates'] if not(pd.isnull(item))==True)


hlalistClassI = {
    "HLAMiner-Class-I": HLAMinerClassI,
    "seq2HLA-Class-I": seq2HLAClassI,
    "Polysolver-Class-I": PolysolverClassI,
    "Athlates-Class-I": AthlatesClassI,
    "Optitype-Class-I": OptitypeClassI
}

hlalist2ClassI = {
    "seq2HLA-Class-I": seq2HLAClassI,
    "Polysolver-Class-I": PolysolverClassI,
    "Athlates-Class-I": AthlatesClassI,
    "Optitype-Class-I": OptitypeClassI
}


#Class II HLA Only
HLAListClassII = pd.read_csv(OutDir+ "/" + SampleID + "-HLA-List-Class-II-Only.csv", header=0,sep=',')
HLAMinerClassII = set(item for item in HLAListClassII['HLAMiner'] if not(pd.isnull(item))==True)
#ClinicalClassII = set(item for item in HLAList['Clinical-HLA'] if not(pd.isnull(item))==True)
seq2HLAClassII = set(item for item in HLAListClassII['Seq2HLA'] if not(pd.isnull(item))==True)
AthlatesClassII = set(item for item in HLAListClassII['Athlates'] if not(pd.isnull(item))==True)


hlalistClassII = {
    "HLAMiner-Class-II": HLAMinerClassII,
    "seq2HLA-Class-II": seq2HLAClassII,
    "Athlates-Class-II": AthlatesClassII
}

#Intrapatient Polysolver Class I HLA Only (DNA-Normal, DNA-Tumor, RNA-Tumor)
PolysolverIntrapatientHLAList = pd.read_csv(OutDir+ "/" + SampleID + "-Polysolver-DNA-Normal-DNA-Tumor-RNA-Tumor.csv", header=0, sep=',')
PolysolverDNANormal = set(item for item in PolysolverIntrapatientHLAList['DNA-Normal'] if not(pd.isnull(item))==True)
PolysolverDNATumor = set(item for item in PolysolverIntrapatientHLAList['DNA-Tumor'] if not(pd.isnull(item))==True)
PolysolverRNATumor = set(item for item in PolysolverIntrapatientHLAList['RNA-Tumor'] if not(pd.isnull(item))==True)

#Polysolver Class I HLA consensus
PolysolverIntrapatienthlalist = {
    "Polysolver-I-DNANormal": PolysolverDNANormal,
    "Polysolver-I-DNATumor": PolysolverDNATumor,
    "Polysolver-I-RNATumor": PolysolverRNATumor
}


#Intrapatient Optitype Class I HLA Only (DNA-Normal, DNA-Tumor, RNA-Tumor)
OptitypeIntrapatientHLAList = pd.read_csv(OutDir+ "/" + SampleID + "-Optitype-DNA-Normal-DNA-Tumor-RNA-Tumor.csv", header=0, sep=',')
OptitypeDNANormal = set(item for item in OptitypeIntrapatientHLAList['DNA-Normal'] if not(pd.isnull(item))==True)
OptitypeDNATumor = set(item for item in OptitypeIntrapatientHLAList['DNA-Tumor'] if not(pd.isnull(item))==True)
OptitypeRNATumor = set(item for item in OptitypeIntrapatientHLAList['RNA-Tumor'] if not(pd.isnull(item))==True)

#Optitype Class I HLA consensus
OptitypeIntrapatienthlalist = {
    "Optitype-I-DNANormal": OptitypeDNANormal,
    "Optitype-I-DNATumor": OptitypeDNATumor,
    "Optitype-I-RNATumor": OptitypeRNATumor
}


#Intrapatient Athlates-Class-I HLA Only (DNA-Normal, DNA-Tumor, RNA-Tumor)
AthlatesClassIIntrapatientHLAList = pd.read_csv(OutDir+ "/" + SampleID + "-Athlates-Class-I-DNA-Normal-DNA-Tumor-RNA-Tumor.csv", header=0, sep=',')
AthlatesClassIDNANormal = set(item for item in AthlatesClassIIntrapatientHLAList['DNA-Normal'] if not(pd.isnull(item))==True)
AthlatesClassIDNATumor = set(item for item in AthlatesClassIIntrapatientHLAList['DNA-Tumor'] if not(pd.isnull(item))==True)
AthlatesClassIRNATumor = set(item for item in AthlatesClassIIntrapatientHLAList['RNA-Tumor'] if not(pd.isnull(item))==True)

#Athlates-Class-I HLA consensus
AthlatesClassIintrapatienthlalist = {
    "Athlates-Class-I-DNANormal": AthlatesClassIDNANormal,
    "Athlates-Class-I-DNATumor": AthlatesClassIDNATumor,
    "Athlates-Class-I-RNATumor": AthlatesClassIRNATumor
}


#Intrapatient Athlates-Class-II HLA Only (DNA-Normal, DNA-Tumor, RNA-Tumor)
AthlatesClassIIIntrapatientHLAList = pd.read_csv(OutDir+ "/" + SampleID + "-Athlates-Class-II-DNA-Normal-DNA-Tumor-RNA-Tumor.csv", header=0, sep=',')
AthlatesClassIIDNANormal = set(item for item in AthlatesClassIIIntrapatientHLAList['DNA-Normal'] if not(pd.isnull(item))==True)
AthlatesClassIIDNATumor = set(item for item in AthlatesClassIIIntrapatientHLAList['DNA-Tumor'] if not(pd.isnull(item))==True)
AthlatesClassIIRNATumor = set(item for item in AthlatesClassIIIntrapatientHLAList['RNA-Tumor'] if not(pd.isnull(item))==True)

#Athlates-Class-I HLA consensus
AthlatesClassIIintrapatienthlalist = {
    "Athlates-Class-II-DNANormal": AthlatesClassIIDNANormal,
    "Athlates-Class-II-DNATumor": AthlatesClassIIDNATumor,
    "Athlates-Class-II-RNATumor": AthlatesClassIIRNATumor
}


#Intrapatient Seq2HLA-Class-I HLA Only (DNA-Normal, DNA-Tumor, RNA-Tumor)
Seq2HLAClassIIntrapatientHLAList = pd.read_csv(OutDir+ "/" + SampleID + "-Seq2HLA-Class-I-DNA-Normal-DNA-Tumor-RNA-Tumor.csv", header=0, sep=',')
Seq2HLAClassIDNANormal = set(item for item in Seq2HLAClassIIntrapatientHLAList['DNA-Normal'] if not(pd.isnull(item))==True)
Seq2HLAClassIDNATumor = set(item for item in Seq2HLAClassIIntrapatientHLAList['DNA-Tumor'] if not(pd.isnull(item))==True)
Seq2HLAClassIRNATumor = set(item for item in Seq2HLAClassIIntrapatientHLAList['RNA-Tumor'] if not(pd.isnull(item))==True)

#Seq2HLA-Class-I HLA consensus
Seq2HLAClassIintrapatienthlalist = {
    "Seq2HLA-Class-I-DNANormal": Seq2HLAClassIDNANormal,
    "Seq2HLA-Class-I-DNATumor": Seq2HLAClassIDNATumor,
    "Seq2HLA-Class-I-RNATumor": Seq2HLAClassIRNATumor
}

#Intrapatient Seq2HLA-Class-II HLA Only (DNA-Normal, DNA-Tumor, RNA-Tumor)
Seq2HLAClassIIIntrapatientHLAList = pd.read_csv(OutDir+ "/" + SampleID + "-Seq2HLA-Class-II-DNA-Normal-DNA-Tumor-RNA-Tumor.csv", header=0, sep=',')
Seq2HLAClassIIDNANormal = set(item for item in Seq2HLAClassIIIntrapatientHLAList['DNA-Normal'] if not(pd.isnull(item))==True)
Seq2HLAClassIIDNATumor = set(item for item in Seq2HLAClassIIIntrapatientHLAList['DNA-Tumor'] if not(pd.isnull(item))==True)
Seq2HLAClassIIRNATumor = set(item for item in Seq2HLAClassIIIntrapatientHLAList['RNA-Tumor'] if not(pd.isnull(item))==True)

#Seq2HLA-Class-I HLA consensus
Seq2HLAClassIIintrapatienthlalist = {
    "Seq2HLA-Class-II-DNANormal": Seq2HLAClassIIDNANormal,
    "Seq2HLA-Class-II-DNATumor": Seq2HLAClassIIDNATumor,
    "Seq2HLA-Class-II-RNATumor": Seq2HLAClassIIRNATumor
}





                                  
venn(hlalistClassI) #Use varible for pptx export VennHSPAO = venn(hlalistClassI)
plt.title(SampleID)
plt.savefig(SampleID + "-HLAMiner-seq2HLA-Polysolver-Optitype-Athlates-ConsensusHLA-2.png")
ImgHSPOAHLACI = SampleID + "-HLAMiner-seq2HLA-Polysolver-Optitype-Athlates-ConsensusHLA-2.png"

venn(hlalist2ClassI)
plt.title(SampleID)
plt.savefig(SampleID + "-seq2HLA-Polysolver-Optitype-Athlates-ConsensusHLA-2.png")
ImgSPOAHLACI = SampleID + "-seq2HLA-Polysolver-Optitype-Athlates-ConsensusHLA-2.png"

venn(hlalistClassII)
plt.title(SampleID)
plt.savefig(SampleID + "-HLAMiner-seq2HLA-Athlates-ConsensusHLA-2.png")
ImgHSAHLACII = SampleID + "-HLAMiner-seq2HLA-Athlates-ConsensusHLA-2.png"

venn(PolysolverIntrapatienthlalist)
plt.title(SampleID)
plt.savefig(SampleID + "-Polysolver-DNA-Normal-DNA-Tumor-RNA-Tumor-2.png")
ImgPolysolverHLACI = SampleID + "-Polysolver-DNA-Normal-DNA-Tumor-RNA-Tumor-2.png"

if not OptitypeIntrapatientHLAList.isnull().values.all(): #venn(OptitypeIntrapatienthlalist) #check if All values NA
    print("&&&&&&&&&&&&&&&&&&&&&&&&&&&&&&&&&&&&&&&&&&&&&&&&&&&&&&&&************^^^^^^^^^^^^^^^^")
    print(OptitypeIntrapatientHLAList.isnull().values.all())
    print(OptitypeIntrapatientHLAList.empty)
    venn(OptitypeIntrapatienthlalist)
    plt.title(SampleID)
    plt.savefig(SampleID + "-Optitype-DNA-Normal-DNA-Tumor-RNA-Tumor-2.png")
    ImgOptitypeHLACI = SampleID + "-Optitype-DNA-Normal-DNA-Tumor-RNA-Tumor-2.png"


venn(AthlatesClassIintrapatienthlalist)
plt.title(SampleID)
plt.savefig(SampleID + "-Athlates-Class-I-DNA-Normal-DNA-Tumor-RNA-Tumor-2.png")
ImgAthlatesHLACI = SampleID + "-Athlates-Class-I-DNA-Normal-DNA-Tumor-RNA-Tumor-2.png"

venn(AthlatesClassIIintrapatienthlalist)
plt.title(SampleID)
plt.savefig(SampleID + "-Athlates-Class-II-DNA-Normal-DNA-Tumor-RNA-Tumor-2.png")
ImgAthlatesHLACII = SampleID + "-Athlates-Class-II-DNA-Normal-DNA-Tumor-RNA-Tumor-2.png"

venn(Seq2HLAClassIintrapatienthlalist)
plt.title(SampleID)
plt.savefig(SampleID + "-Seq2HLA-Class-I-DNA-Normal-DNA-Tumor-RNA-Tumor-2.png")
ImgSeq2HLACI = SampleID + "-Seq2HLA-Class-I-DNA-Normal-DNA-Tumor-RNA-Tumor-2.png"

venn(Seq2HLAClassIIintrapatienthlalist)
plt.title(SampleID)
plt.savefig(SampleID + "-Seq2HLA-Class-II-DNA-Normal-DNA-Tumor-RNA-Tumor-2.png")
ImgSeq2HLACII = SampleID + "-Seq2HLA-Class-II-DNA-Normal-DNA-Tumor-RNA-Tumor-2.png"

#venn(intrapatienthlalist)

print("Common HLA (HLAMinerClassI, seq2HLAClassI, PolysolverClassI, OptitypeClassI, AthlatesClassI):", set.intersection(HLAMinerClassI, seq2HLAClassI, PolysolverClassI, OptitypeClassI, AthlatesClassI))
print("Common HLA (HLAMinerClassI, seq2HLAClassI, PolysolverClassI, AthlatesClassI):", set.intersection(HLAMinerClassI, seq2HLAClassI, PolysolverClassI, AthlatesClassI))
print("Common HLA (seq2HLAClassI, PolysolverClassI, AthlatesClassI):", set.intersection(seq2HLAClassI, PolysolverClassI, AthlatesClassI))
print("Common HLA (seq2HLAClassI, PolysolverClassI, OptitypeClassI, AthlatesClassI):", set.intersection(seq2HLAClassI, PolysolverClassI, OptitypeClassI, AthlatesClassI))
','.join(HLAMinerClassI&seq2HLAClassI&PolysolverClassI&OptitypeClassI&AthlatesClassI)
print("Common HLA (seq2HLAClassI, PolysolverClassI, OptitypeClassI):", set.intersection(seq2HLAClassI, PolysolverClassI, OptitypeClassI))
','.join(HLAMinerClassI&seq2HLAClassI&PolysolverClassI&OptitypeClassI)
print("Common HLA (PolysolverClassI, OptitypeClassI, AthlatesClassI):", set.intersection(PolysolverClassI, OptitypeClassI, AthlatesClassI))
','.join(HLAMinerClassI&seq2HLAClassI&PolysolverClassI&OptitypeClassI&AthlatesClassI)
#' '.join(seq2HLA&Polysolver&Optitype)
print("Common HLA (HLAMinerClassII, seq2HLAClassII, AthlatesClassII):", set.intersection(HLAMinerClassII, seq2HLAClassII, AthlatesClassII))
print("Common HLA (seq2HLAClassII, AthlatesClassII):", set.intersection(seq2HLAClassII, AthlatesClassII))
print("Common HLA (HLAMinerClassII, AthlatesClassII):", set.intersection(HLAMinerClassII, AthlatesClassII))
print("Common HLA (HLAMinerClassII, seq2HLAClassII):", set.intersection(HLAMinerClassII, seq2HLAClassII))





#Create pptx report
from datetime import date
from pptx import Presentation
from pptx.util import Inches, Pt


DateToday = date.today().strftime("%m/%d/%Y")

HLAppt = Presentation()
HLAppt.slide_width=Inches(13.333)
HLAppt.slide_height=Inches(7.5)
TileSlide = HLAppt.slide_layouts[0]
HLASlide = HLAppt.slide_layouts[5]


""" Ref for slide types: 
0 ->  title and subtitle
1 ->  title and content
2 ->  section header
3 ->  two content
4 ->  Comparison
5 ->  Title only 
6 ->  Blank
7 ->  Content with caption
8 ->  Pic with caption
"""
  
#Add title slide
slide01 = HLAppt.slides.add_slide(TileSlide)
slide01.shapes.title.width = Inches(11)
slide01.shapes.title.height = Inches(2)
slide01.shapes.title.left = Inches(1)
slide01.shapes.title.top = Inches(2)
slide01.shapes.title.text = "Consensus HLA Analysis"
slide01.placeholders[1].width = Inches(9)
slide01.placeholders[1].height = Inches(2)
slide01.placeholders[1].left = Inches(2)
slide01.placeholders[1].top = Inches(4)
slide01.placeholders[1].text = SampleID 

AllTooldf = {1 : ["AthlatesClassIIntrapatientHLAList", "slide02", "Athlates Class I", "ImgAthlatesHLACI"], 2 : ["PolysolverIntrapatientHLAList", "slide03", "Polysolver Class I", "ImgPolysolverHLACI"], 3 : ["OptitypeIntrapatientHLAList", "slide04", "Optitype Class I", "ImgOptitypeHLACI"], 4 : ["Seq2HLAClassIIntrapatientHLAList", "slide05", "Seq2HLA Class I", "ImgSeq2HLACI"], 5 : ["AthlatesClassIIIntrapatientHLAList", "slide06", "Athlates Class II", "ImgAthlatesHLACII"], 6 : ["Seq2HLAClassIIIntrapatientHLAList", "slide07", "Seq2HLA Class II", "ImgSeq2HLACII"]}


for M, N in AllTooldf.items():
    print(M, N)
    Vars = vars()
    DataframE = Vars[str(N[0])] #AthlatesClassIIntrapatientHLAList #evaluate string name to data frame
    try:
        VennImg = Vars[str(N[3])]
    except:
        VennImg = None
    N[1] = HLAppt.slides.add_slide(HLASlide)
    N[1].shapes.title.width = Inches(11)
    N[1].shapes.title.height = Inches(1)
    N[1].shapes.title.left = Inches(1)
    N[1].shapes.title.text = SampleID
    text_box1 = N[1].shapes.add_textbox(Inches(11), Inches(1), Inches(1), Inches(0.5)) #(left, top, width, height)
    tb1=text_box1.text_frame
    tb1.text = DateToday
    text_box2 = N[1].shapes.add_textbox(Inches(1), Inches(1), Inches(2), Inches(0.5)) #(left, top, width, height)
    tb2=text_box2.text_frame
    tb2.text = N[2] #Athlates Class I"
    print("type of DataFramE: ", type(DataframE))
    Rows = len(DataframE)
    Columns = len(DataframE.columns)
    try:
        N[1].shapes.add_picture(VennImg, Inches(6.5), Inches(1), Inches(6.5), Inches(6.5)) #ImgAthlatesHLACI
    except:
        VennImg = None
    Table = N[1].shapes.add_table(Rows + 1, Columns, Inches(0.5), Inches(1.5), Inches(6.5), Inches(Rows * 0.5))

    for Name in range(len(DataframE.columns)):
        Table.table.cell(0, Name).text = str(DataframE.columns.values[Name])

    print(DataframE)
    print(Rows)
    print(Columns)
    for i in range(Rows):
        for j in range(Columns):
            print(i, j)
            print(DataframE.values[i, j])
            Table.table.cell(i+1, j).text = str(DataframE.values[i, j])
#    HLAppt.save(SampleID + "ConsensusHLA-Analysis.pptx")


#table for All Class I consensus HLA
DataframE = HLAListClassI
slide08 = HLAppt.slides.add_slide(HLASlide)
slide08.shapes.title.width = Inches(11)
slide08.shapes.title.height = Inches(1)
slide08.shapes.title.left = Inches(1)
slide08.shapes.title.text = SampleID
text_box1 = slide08.shapes.add_textbox(Inches(11), Inches(1), Inches(1), Inches(0.5)) #(left, top, width, height)
tb1=text_box1.text_frame
tb1.text = DateToday
text_box2 = slide08.shapes.add_textbox(Inches(1), Inches(1), Inches(3), Inches(0.5)) #(left, top, width, height)
tb2=text_box2.text_frame
tb2.text = "All Class I consensus HLA"
Rows = len(DataframE)
Columns = len(DataframE.columns)
#try:
#    slide08.shapes.add_picture(VennImg, Inches(6.5), Inches(1), Inches(6.5), Inches(6.5)) #ImgAthlatesHLACI
#except:
#    VennImg = None
Table = slide08.shapes.add_table(Rows + 1, Columns, Inches(0.5), Inches(1.5), Inches(11), Inches(Rows * 0.5))

for Name in range(len(DataframE.columns)):
    Table.table.cell(0, Name).text = str(DataframE.columns.values[Name])

print(DataframE)
print(Rows)
print(Columns)
for i in range(Rows):
    for j in range(Columns):
        print(i, j)
        print(DataframE.values[i, j])
        Table.table.cell(i+1, j).text = str(DataframE.values[i, j])


#Figures for All Class I consensus HLA
DataframE = HLAListClassI
slide09 = HLAppt.slides.add_slide(HLASlide)
slide09.shapes.title.width = Inches(11)
slide09.shapes.title.height = Inches(1)
slide09.shapes.title.left = Inches(1)
slide09.shapes.title.text = SampleID
text_box1 = slide09.shapes.add_textbox(Inches(11), Inches(1), Inches(1), Inches(0.5)) #(left, top, width, height)
tb1=text_box1.text_frame
tb1.text = DateToday
text_box2 = slide09.shapes.add_textbox(Inches(1), Inches(1), Inches(3), Inches(0.5)) #(left, top, width, height)
tb2=text_box2.text_frame
tb2.text = "All Class I consensus HLA"
Rows = len(DataframE)
Columns = len(DataframE.columns)
try:
    slide09.shapes.add_picture(ImgHSPOAHLACI, Inches(0.5), Inches(1), Inches(6.5), Inches(6.5)) #ImgAthlatesHLACI
except:
    ImgHSPOAHLACI = None
try:
    slide09.shapes.add_picture(ImgSPOAHLACI, Inches(6.5), Inches(1), Inches(6.5), Inches(6.5)) #ImgAthlatesHLACI
except:
    ImgSPOAHLACI = None
text_box3 = slide09.shapes.add_textbox(Inches(0.5), Inches(6.75), Inches(12), Inches(1)) #(left, top, width, height)
tb3=text_box3.text_frame
tb3.text = "All Class I consensus HLA: " + str(list(ConsensusClassIHLA)).replace("'", "")



#table for All Class II consensus HLA
DataframE = HLAListClassII
slide10 = HLAppt.slides.add_slide(HLASlide)
slide10.shapes.title.width = Inches(11)
slide10.shapes.title.height = Inches(1)
slide10.shapes.title.left = Inches(1)
slide10.shapes.title.text = SampleID
text_box1 = slide10.shapes.add_textbox(Inches(11), Inches(1), Inches(1), Inches(0.5)) #(left, top, width, height)
tb1=text_box1.text_frame
tb1.text = DateToday
text_box2 = slide10.shapes.add_textbox(Inches(1), Inches(1), Inches(3), Inches(0.5)) #(left, top, width, height)
tb2=text_box2.text_frame
tb2.text = "All Class II consensus HLA"
Rows = len(DataframE)
Columns = len(DataframE.columns)
#try:
#    slide08.shapes.add_picture(VennImg, Inches(6.5), Inches(1), Inches(6.5), Inches(6.5)) #ImgAthlatesHLACI
#except:
#    VennImg = None
Table = slide10.shapes.add_table(Rows + 1, Columns, Inches(0.5), Inches(1.5), Inches(11), Inches(Rows * 0.5))

for Name in range(len(DataframE.columns)):
    Table.table.cell(0, Name).text = str(DataframE.columns.values[Name])

print(DataframE)
print(Rows)
print(Columns)
for i in range(Rows):
    for j in range(Columns):
        print(i, j)
        print(DataframE.values[i, j])
        Table.table.cell(i+1, j).text = str(DataframE.values[i, j])


#Figures for All Class II consensus HLA
slide11 = HLAppt.slides.add_slide(HLASlide)
slide11.shapes.title.width = Inches(11)
slide11.shapes.title.height = Inches(1)
slide11.shapes.title.left = Inches(1)
slide11.shapes.title.text = SampleID
text_box1 = slide11.shapes.add_textbox(Inches(11), Inches(1), Inches(1), Inches(0.5)) #(left, top, width, height)
tb1=text_box1.text_frame
tb1.text = DateToday
text_box2 = slide11.shapes.add_textbox(Inches(1), Inches(1), Inches(3), Inches(0.5)) #(left, top, width, height)
tb2=text_box2.text_frame
tb2.text = "All Class II consensus HLA"
try:
    slide11.shapes.add_picture(ImgHSAHLACII, Inches(0.5), Inches(1), Inches(6.5), Inches(6.5)) #ImgAthlatesHLACI
except:
    ImgHSPOAHLACII = None
text_box3 = slide11.shapes.add_textbox(Inches(0.5), Inches(6.75), Inches(12), Inches(1)) #(left, top, width, height)
tb3=text_box3.text_frame
tb3.text = "All Class II consensus HLA: " + str(list(ConsensusClassIIHLA)).replace("'", "")

    


#print(Dataframe.columns)

# Saving file
HLAppt.save(SampleID + "ConsensusHLA-Analysis.pptx")

print("ppt report created")


